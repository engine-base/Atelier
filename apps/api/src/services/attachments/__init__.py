"""GAP-161: 添付資料を AI が実際に参照できる形にする。

経営者指摘 (2026-08-19):
  「デザインモックも、このテンプレもだけど画像や PDF やファイルやエクセルを
   アップロードしてそれを参考にすることがチャットでできていないけどどうして？？」

実バグだった: chat_messages.attachments に保存・画面表示はしていたが、
**LLM への引き渡しが一切無かった** (system prompt にもユーザー文にも入らない)。

設計:
  - PDF / Excel / Word / PowerPoint / CSV / テキストは **サーバー側で決定的に
    テキスト抽出** し、プロンプトへ注入する。LLM を使わないので追加費用ゼロで、
    かつ Bridge (本人サブスク) / サブスク / API のどの実行経路でも同じに効く。
  - 画像は文字起こしを推測しない。ファイル名と形式を伝えたうえで、
    Bridge 経路ではローカル作業場へ実体を配って Claude Code に直接読ませる
    (chat_relay の workspace seed 拡張 — 同 GAP)。
  - 抽出できない形式・壊れたファイルは **正直にその旨を書く** (無言で落とさない)。
"""

from __future__ import annotations

import csv
import io
from dataclasses import dataclass

# 1 添付あたりの注入上限 (長大な資料でプロンプトを溢れさせない)
MAX_CHARS_PER_ATTACHMENT = 12_000
# 全添付合計の上限
MAX_CHARS_TOTAL = 30_000

IMAGE_MIME = frozenset({"image/png", "image/jpeg", "image/webp", "image/gif"})

_SHEET_MAX_ROWS = 200
_SHEET_MAX_COLS = 40


@dataclass(frozen=True)
class ExtractedAttachment:
    """添付 1 件の抽出結果。text が空でも kind/note で状況を正直に伝える。"""

    file_name: str
    mime_type: str
    kind: str  # "text" | "image" | "unsupported" | "error"
    text: str
    note: str = ""


def _decode(data: bytes) -> str:
    for enc in ("utf-8", "cp932", "utf-16"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _clip(text: str) -> str:
    t = text.strip()
    if len(t) <= MAX_CHARS_PER_ATTACHMENT:
        return t
    return t[:MAX_CHARS_PER_ATTACHMENT] + "\n…(以降は長いため省略)"


def _extract_pdf(data: bytes) -> tuple[str, str]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        body = (page.extract_text() or "").strip()
        if body:
            pages.append(f"[p.{i}]\n{body}")
    if not pages:
        return "", "PDF からテキストを抽出できませんでした (画像 PDF の可能性があります)"
    return "\n\n".join(pages), f"{len(reader.pages)} ページ"


def _extract_xlsx(data: bytes) -> tuple[str, str]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    chunks: list[str] = []
    for ws in wb.worksheets:
        rows: list[str] = []
        for r_i, row in enumerate(ws.iter_rows(values_only=True)):
            if r_i >= _SHEET_MAX_ROWS:
                rows.append("…(以降の行は省略)")
                break
            cells = ["" if c is None else str(c) for c in row[:_SHEET_MAX_COLS]]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells).rstrip(" |"))
        if rows:
            chunks.append(f"## シート: {ws.title}\n" + "\n".join(rows))
    wb.close()
    if not chunks:
        return "", "空のブックでした"
    return "\n\n".join(chunks), f"{len(wb.worksheets)} シート"


def _extract_docx(data: bytes) -> tuple[str, str]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            line = " | ".join(c.text.strip() for c in row.cells)
            if line.strip(" |"):
                parts.append(line)
    if not parts:
        return "", "本文が空でした"
    return "\n".join(parts), f"{len(doc.paragraphs)} 段落"


def _extract_pptx(data: bytes) -> tuple[str, str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [
            str(shape.text).strip()  # pyright: ignore[reportAttributeAccessIssue]
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and str(shape.text).strip()  # pyright: ignore[reportAttributeAccessIssue]
        ]
        if lines:
            slides.append(f"[スライド {i}]\n" + "\n".join(lines))
    if not slides:
        return "", "テキストのあるスライドがありませんでした"
    return "\n\n".join(slides), f"{len(prs.slides)} スライド"


def _extract_csv(data: bytes) -> tuple[str, str]:
    text = _decode(data)
    reader = csv.reader(io.StringIO(text))
    rows = []
    for i, row in enumerate(reader):
        if i >= _SHEET_MAX_ROWS:
            rows.append("…(以降の行は省略)")
            break
        rows.append(" | ".join(row[:_SHEET_MAX_COLS]))
    return "\n".join(rows), f"{len(rows)} 行"


def extract_attachment(*, file_name: str, mime_type: str, data: bytes) -> ExtractedAttachment:
    """添付 1 件を AI が読めるテキストにする (LLM 不使用・決定的)。"""
    mime = (mime_type or "").split(";")[0].strip().lower()
    name = file_name or "(名称不明)"
    lower = name.lower()

    if mime in IMAGE_MIME:
        return ExtractedAttachment(
            file_name=name,
            mime_type=mime,
            kind="image",
            text="",
            note="画像 (内容の推測はしない — 実行経路が対応していれば実物を参照する)",
        )

    try:
        if mime == "application/pdf" or lower.endswith(".pdf"):
            body, note = _extract_pdf(data)
        elif mime in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ) or lower.endswith((".xlsx", ".xlsm")):
            body, note = _extract_xlsx(data)
        elif (
            mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or lower.endswith(".docx")
        ):
            body, note = _extract_docx(data)
        elif (
            mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or lower.endswith(".pptx")
        ):
            body, note = _extract_pptx(data)
        elif mime == "text/csv" or lower.endswith(".csv"):
            body, note = _extract_csv(data)
        elif mime.startswith("text/") or lower.endswith((".txt", ".md", ".json", ".yaml", ".yml")):
            body, note = _decode(data), ""
        else:
            return ExtractedAttachment(
                file_name=name,
                mime_type=mime,
                kind="unsupported",
                text="",
                note=f"この形式 ({mime or '不明'}) はテキスト抽出に未対応です",
            )
    except Exception as exc:  # 壊れたファイル等 — 無言で捨てず理由を残す
        return ExtractedAttachment(
            file_name=name,
            mime_type=mime,
            kind="error",
            text="",
            note=f"読み取りに失敗しました: {type(exc).__name__}",
        )

    if not body.strip():
        return ExtractedAttachment(
            file_name=name,
            mime_type=mime,
            kind="error",
            text="",
            note=note or "内容を取り出せませんでした",
        )
    return ExtractedAttachment(
        file_name=name, mime_type=mime, kind="text", text=_clip(body), note=note
    )


def render_attachments_block(items: list[ExtractedAttachment]) -> str:
    """抽出結果を system prompt 用のブロックにする。空なら ""。"""
    if not items:
        return ""
    lines = [
        "# 添付資料 (ユーザーがこの会話にアップロードしたファイルの中身)",
        "以下はユーザーが参考として渡した資料の実際の内容です。"
        "依頼に関係する部分は必ず踏まえて回答・作成すること。"
        "内容が取り出せなかったものは、推測で補わずその旨を伝えること。",
    ]
    used = 0
    for it in items:
        head = f"\n## {it.file_name}（{it.mime_type or '形式不明'}{f' / {it.note}' if it.note else ''}）"
        if it.kind == "text":
            remain = MAX_CHARS_TOTAL - used
            if remain <= 0:
                lines.append(f"{head}\n(全体の分量上限に達したため本文は省略)")
                continue
            body = it.text[:remain]
            used += len(body)
            lines.append(f"{head}\n{body}")
        elif it.kind == "image":
            lines.append(f"{head}\n(画像。内容を推測しないこと)")
        else:
            lines.append(f"{head}\n(取り込めませんでした: {it.note})")
    return "\n".join(lines)


async def extract_stored_attachments(
    records: list[dict[str, object]], *, max_files: int = 5
) -> list[ExtractedAttachment]:
    """storage 上の添付/参考資料をまとめて取得・抽出する (GAP-161)。

    取得失敗は例外にせず kind="error" として残す — 会話や生成自体は止めない。
    """
    import httpx

    from src.storage_signing import create_signed_download_url

    out: list[ExtractedAttachment] = []
    seen: set[str] = set()
    for rec in records:
        path = str(rec.get("storage_path") or "")
        if not path or path in seen:
            continue
        seen.add(path)
        name = str(rec.get("file_name") or "(名称不明)")
        mime = str(rec.get("mime_type") or "")
        try:
            url = await create_signed_download_url(path)
            async with httpx.AsyncClient(timeout=20.0) as client:
                res = await client.get(url)
            if res.status_code >= 400:
                raise RuntimeError(f"status {res.status_code}")
            data = res.content
        except Exception as exc:
            out.append(
                ExtractedAttachment(
                    file_name=name,
                    mime_type=mime,
                    kind="error",
                    text="",
                    note=f"ファイル取得に失敗しました ({type(exc).__name__})",
                )
            )
            continue
        out.append(extract_attachment(file_name=name, mime_type=mime, data=data))
        if len(out) >= max_files:
            break
    return out


def render_reference_block(items: list[ExtractedAttachment]) -> str:
    """スタジオ (モック/デザインテンプレ) 用の参考資料ブロック (GAP-161)。

    「これを参考にしろ。ただし推測で埋めるな」を明示する点は chat と同じだが、
    デザイン作業向けに参照の仕方を書き分ける。
    """
    if not items:
        return ""
    lines = [
        "# 参考資料 (ユーザーがこの作業のためにアップロードした資料)",
        "レイアウト・配色・項目立て・文言のトーンなど、依頼に関係する点は"
        "この資料に合わせること。画像は内容を推測せず、読み取れないものは"
        "無理に反映しないこと。",
    ]
    used = 0
    for it in items:
        head = f"\n## {it.file_name}（{it.mime_type or '形式不明'}{f' / {it.note}' if it.note else ''}）"
        if it.kind == "text":
            remain = MAX_CHARS_TOTAL - used
            if remain <= 0:
                lines.append(f"{head}\n(分量上限のため省略)")
                continue
            body = it.text[:remain]
            used += len(body)
            lines.append(f"{head}\n{body}")
        elif it.kind == "image":
            lines.append(f"{head}\n(画像。内容を推測しないこと)")
        else:
            lines.append(f"{head}\n(取り込めませんでした: {it.note})")
    return "\n".join(lines)
