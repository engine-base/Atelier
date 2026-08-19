"""GAP-161: 添付資料のテキスト抽出 (LLM 不使用・決定的) のテスト。"""

from __future__ import annotations

import io

from src.services.attachments import (
    ExtractedAttachment,
    extract_attachment,
    render_attachments_block,
)


def _xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "見積"
    ws.append(["項目", "数量", "単価"])
    ws.append(["設計", 1, 300000])
    ws.append(["実装", 2, 500000])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _docx_bytes() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph("要件定義メモ")
    doc.add_paragraph("会員登録は SNS ログインを含む")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "提案の骨子"  # pyright: ignore[reportAttributeAccessIssue, reportOptionalMemberAccess]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pdf_bytes() -> bytes:
    from reportlab.pdfgen import canvas

    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.drawString(100, 700, "Estimate total 800000 JPY")
    c.showPage()
    c.save()
    return buf.getvalue()


class TestExtract:
    def test_excel_becomes_readable_rows(self) -> None:
        got = extract_attachment(
            file_name="見積.xlsx",
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            data=_xlsx_bytes(),
        )
        assert got.kind == "text"
        assert "シート: 見積" in got.text
        assert "設計 | 1 | 300000" in got.text

    def test_word_and_powerpoint(self) -> None:
        w = extract_attachment(
            file_name="要件.docx",
            mime_type=("application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
            data=_docx_bytes(),
        )
        assert w.kind == "text" and "SNS ログイン" in w.text
        p = extract_attachment(
            file_name="提案.pptx",
            mime_type=("application/vnd.openxmlformats-officedocument.presentationml.presentation"),
            data=_pptx_bytes(),
        )
        assert p.kind == "text" and "提案の骨子" in p.text

    def test_pdf_and_csv(self) -> None:
        pdf = extract_attachment(
            file_name="estimate.pdf", mime_type="application/pdf", data=_pdf_bytes()
        )
        assert pdf.kind == "text" and "800000" in pdf.text
        csv_ = extract_attachment(
            file_name="rows.csv",
            mime_type="text/csv",
            data="名前,金額\n設計,300000\n".encode(),
        )
        assert csv_.kind == "text" and "設計 | 300000" in csv_.text

    def test_image_is_not_guessed_and_broken_file_is_honest(self) -> None:
        img = extract_attachment(
            file_name="ui.png", mime_type="image/png", data=b"\x89PNG\r\n\x1a\n"
        )
        assert img.kind == "image" and img.text == ""
        assert "推測" in img.note
        broken = extract_attachment(
            file_name="broken.pdf", mime_type="application/pdf", data=b"not a pdf at all"
        )
        assert broken.kind == "error" and broken.text == ""
        unsupported = extract_attachment(
            file_name="movie.mp4", mime_type="video/mp4", data=b"\x00\x00"
        )
        assert unsupported.kind == "unsupported"
        assert "未対応" in unsupported.note

    def test_block_tells_ai_to_use_it_and_not_to_invent(self) -> None:
        block = render_attachments_block(
            [
                ExtractedAttachment(
                    file_name="見積.xlsx",
                    mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    kind="text",
                    text="設計 | 300000",
                    note="1 シート",
                ),
                ExtractedAttachment(
                    file_name="ui.png",
                    mime_type="image/png",
                    kind="image",
                    text="",
                    note="画像",
                ),
            ]
        )
        assert "# 添付資料" in block
        assert "見積.xlsx" in block and "設計 | 300000" in block
        assert "推測で補わず" in block
        assert "画像。内容を推測しないこと" in block
        assert render_attachments_block([]) == ""
